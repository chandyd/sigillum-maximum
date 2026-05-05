#!/usr/bin/env python3
"""
Generate ALL pitch PDFs for the entire RCP / Eleanor Lian ecosystem:
1. RCP Master Pitch (all IPs together)
2. Master Pitch for each IP (Sigillum, Clown, Bein)
3. Derivative pitches for each IP (anime, film, game, merch, manga, board game, experiences)

All saved to their respective site folders under sigillum-maximum-site/derivative-pitches/
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

GOLD = RGBColor(0xD4, 0xA7, 0x47)
DARK = RGBColor(0x0A, 0x0A, 0x0F)
CREAM = RGBColor(0xF0, 0xE8, 0xD8)
LIGHT = RGBColor(0xE0, 0xDC, 0xD5)
MUTED = RGBColor(0xA0, 0x98, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gold_bar(slide, top=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def add_title(slide, text, left=0.8, top=0.5, size=28):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8.4), Inches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = CREAM
    p.font.bold = True
    p.font.name = 'Calibri'

def add_subtitle(slide, text, top=1.2):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.4), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = MUTED
    p.font.name = 'Calibri'

def add_body(slide, items, left=0.8, top=1.7, width=8.4, size=13, color=LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(5)
    return txBox

def add_cta(slide, text="Interested? Contact: info@robertocalvoproductions.com"):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = GOLD
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

def make_pitch(path, slides):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_gold_bar(slide)
        add_title(slide, data['title'], size=data.get('title_size', 28))
        add_subtitle(slide, data.get('subtitle', ''))
        if 'body' in data:
            add_body(slide, data['body'])
        if data.get('has_cta'):
            add_cta(slide)
    prs.save(path)
    return os.path.getsize(path)

outdir = "/root/.openclaw/workspace/sigillum-maximum-site/derivative-pitches"
os.makedirs(outdir, exist_ok=True)

# ===========================
# RCP MASTER PITCH (all IPs)
# ===========================
make_pitch(f"{outdir}/RCP_Eleanor_Lian_Master_Pitch.pdf", [
    {
        'title': 'RCP Productions — Eleanor Lian IP Portfolio',
        'subtitle': 'Three original fantasy universes. Endless adaptation possibilities.',
        'body': [
            'ROBERTO CALVO PRODUCTIONS LTD presents the complete portfolio of Eleanor Lian,',
            'an author whose work spans epic high fantasy, dark fairy tales, and emotional color-magic stories.',
            '',
            'Each IP below is fully developed with:',
            '• Complete manuscript / story bible',
            '• Character design sheets and world concept art',
            '• Defined transmedia roadmap across 6+ commercial derivatives',
            '• Ready for licensing, co-production, or outright acquisition',
            '',
            '---',
            '',
            'IP 1 — SIGILLUM MAXIMUM: The Seal of the Seven Secrets',
            'Genre: Epic High Fantasy | Format: Anime / Film / Game / Manga / More',
            'Status: Complete manuscript (~100K words), 55+ concept images, landing page live',
            '',
            'IP 2 — THE CLOWN AND THE ICE CASTLE',
            'Genre: Dark Fantasy / Mystery | Format: Film / Series / Graphic Novel',
            'Status: Complete manuscript (85K words, 35 chapters), character art available',
            '',
            'IP 3 — BEIN AND THE WORLD OF THE COLORS',
            'Genre: Emotional Fantasy / Coming-of-Age | Format: Animation / Children\'s Book / Game',
            'Status: Complete manuscript, 23+ concept images, full character roster',
            '',
            '---',
            'Contact RCP for full literary rights, adaptation rights, and co-production inquiries.',
        ],
        'has_cta': True,
    }
])

# ===========================
# SIGILLUM MAXIMUM — MASTER PITCH (already exists, but make a PDF version too)
# ===========================
make_pitch(f"{outdir}/Sigillum_Maximum_Master_Pitch.pdf", [
    {'title': 'SIGILLUM MAXIMUM', 'subtitle': 'The Seal of the Seven Secrets — IP Overview', 'body': [
        'Logline: A princess must unlock seven ancient secrets across seven magical realms',
        'to save her kingdom — and discover the truth about who she really is.',
        '',
        'Author: Eleanor Lian | Publisher: RCP Productions Ltd',
        'Status: Complete manuscript, concept art, landing page at sigillummaximum.com',
        '',
        'STORY SUMMARY',
        'Princess Aileen of Nuvolandia witnesses her kingdom shattered by the mysterious',
        'Black Knight Dorcha. Her mother turned to stone, her father captured, she escapes',
        'through a hidden passage to seek the Seven Secrets of the Sigillum Maximum —',
        'ancient seals that hold the power to reshape reality itself.',
        '',
        'Alongside the Knight of Golden Light Aeltiàfisar, the shape-shifting Grogher,',
        'the cursed princess Majory, and the winged leonine Sidae, Aileen journeys through',
        'seven kingdoms: Nuvolandia, Elf, Gnome, Merfolk, Dwarf, Storm, and the legendary',
        'Realm of the Great Light.',
        '',
        'THEME: Every secret has a cost. Every choice has a consequence.',
        '',
        'TRANSMEDIA READINESS: This IP is designed for simultaneous development across:',
        'Anime Series, Film Trilogy, Video Game, Manga, Board Game, Merchandising, Experiences.',
        '',
        'RIGHTS AVAILABLE: Worldwide adaptation and licensing rights.',
    ], 'has_cta': True},
])

# ===========================
# THE CLOWN AND THE ICE CASTLE — ALL PITCHES
# ===========================
make_pitch(f"{outdir}/Clown_Ice_Castle_Master_Pitch.pdf", [
    {'title': 'THE CLOWN AND THE ICE CASTLE', 'subtitle': 'A Dark Fantasy Tale by Eleanor Lian', 'body': [
        'Genre: Dark Fantasy / Mystery / Psychological Drama',
        'Format: Novel (complete, ~85K words, 35 chapters)',
        'Status: Complete manuscript, character art, 4 concept images',
        '',
        'LOGLINE: A young woman and a mysterious clown venture into an enchanted ice castle',
        'that holds the secrets of a forgotten kingdom — and the truth about herself.',
        '',
        'STORY: In a land haunted by eternal winter, Thalìa discovers a hidden path to an ice',
        'castle that appears only under the full moon. With George, a quiet companion, she',
        'enters a world where the boundary between reality and nightmare dissolves. The',
        'castle\'s clown — both guide and trickster — leads her through frozen halls where',
        'each room reveals a buried memory, a forgotten kingdom, and a choice that will',
        'melt or shatter everything she believes.',
        '',
        'THEMES: Memory, identity, the masks we wear, the cost of truth.',
        '',
        'KEY CHARACTERS: Thalìa (protagonist), George (childhood friend / adult companion),',
        'The Clown (enigmatic guide / trickster), Ice Castle (sentient setting / antagonist).',
        '',
        'TONE: Labyrinth meets The NeverEnding Story meets Paprika.',
        '',
        'RIGHTS AVAILABLE: Film, Series, Graphic Novel, Game, Merchandising.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Clown_Ice_Castle_Film.pdf", [
    {'title': 'THE CLOWN AND THE ICE CASTLE — Film', 'subtitle': 'A dark fantasy feature film', 'body': [
        'Format: Feature Film (100-120 min) — Live-Action or High-End Animation',
        'Target: Young Adult / Adult fantasy audiences (PG-13)',
        '',
        'PITCH: A psychological dark fantasy in the tradition of Pan\'s Labyrinth and',
        'The Fall. The Ice Castle itself is a character — a labyrinthine palace that',
        'shifts its halls and remembers everything. Each room challenges the protagonist',
        'with a moral choice disguised as a puzzle.',
        '',
        'VISUAL STYLE: Gothic fantasy architecture meets crystalline ice palaces.',
        'Contrast between the warm, organic human world and the cold, geometric',
        'perfection of the ice castle.',
        '',
        'STRUCTURE: Three acts — (1) Discovery & Entry, (2) Descent into Memory,',
        '(3) The Truth & The Choice.',
        '',
        'Comparable: Pan\'s Labyrinth, The Fall, Crimson Peak (visually),',
        'The Shape of Water (emotionally).',
        '',
        'Available: Film adaptation rights worldwide.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Clown_Ice_Castle_Series.pdf", [
    {'title': 'THE CLOWN AND THE ICE CASTLE — Limited Series', 'subtitle': 'A dark fantasy series for streaming', 'body': [
        'Format: Limited Series, 6-8 episodes × 45-50 min',
        'Target: Adult fantasy / mystery audiences (TV-MA / 16+)',
        'Platform: Netflix, HBO, Apple TV+, Amazon Prime',
        '',
        'PITCH: Each episode explores one "room" or level of the Ice Castle,',
        'uncovering a layer of memory and mystery. The series format allows for',
        'deeper character work and more complex narrative layering than a film.',
        '',
        'EPISODE STRUCTURE:',
        '• Ep 1-2: Discovery — Thalìa finds the path, enters the castle',
        '• Ep 3-5: Descent — Each floor reveals a different kingdom memory',
        '• Ep 6: The Clown\'s Confession — the trickster\'s true nature revealed',
        '• Ep 7-8: The Throne Room and the Final Choice',
        '',
        'TONE: Gothic mystery with emotional depth. Comparable to The OA meets',
        'The Haunting of Hill House meets Dark.',
        '',
        'Available: Series adaptation rights worldwide.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Clown_Ice_Castle_Graphic_Novel.pdf", [
    {'title': 'THE CLOWN AND THE ICE CASTLE — Graphic Novel', 'subtitle': 'A dark fantasy graphic novel adaptation', 'body': [
        'Format: Graphic Novel, 3-4 volumes, ~200 pages each',
        'Target: Manga / graphic novel readers (16+)',
        '',
        'PITCH: The Ice Castle\'s visual richness makes it a perfect fit for the',
        'graphic novel medium. Each volume could explore different floors of the',
        'castle with distinct art styles matching the emotional tone.',
        '',
        'VOLUMES:',
        '• Vol 1: The Frozen Threshold — entering the castle',
        '• Vol 2: The Hall of Mirrors — confronting memories',
        '• Vol 3: The Heart of Ice — the clown\'s domain',
        '• Vol 4: The Thaw — resolution and return',
        '',
        'ART STYLE: Gothic manga meets European bande dessinée.',
        'Monochrome with strategic color for magical elements.',
        '',
        'Comparable: The Sandman, Monstress, Through the Woods.',
        '',
        'Available: Comic / graphic novel publishing rights worldwide.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Clown_Ice_Castle_Game.pdf", [
    {'title': 'THE CLOWN AND THE ICE CASTLE — Video Game', 'subtitle': 'Psychological horror puzzle-adventure', 'body': [
        'Format: Puzzle-Adventure / Psychological Horror',
        'Platform: PC, Consoles, Mobile (tablet-friendly)',
        'Target: Fans of Portal, The Witness, Little Nightmares',
        '',
        'PITCH: A first-person puzzle game set inside the Ice Castle.',
        'Each floor is a self-contained puzzle box with a story layer.',
        'The Clown appears at key moments — sometimes guide, sometimes threat.',
        '',
        'CORE MECHANICS:',
        '• Ice manipulation — freeze water to create platforms, melt barriers',
        '• Memory puzzles — reconstruct fragmented memories to open doors',
        '• Mirror navigation — one mirror shows the past, one shows the truth',
        '• The Clown\'s game — the trickster challenges you to moral choices',
        '',
        'VISUAL: Gothic ice architecture, dynamic lighting, reflective surfaces.',
        'Sound design is critical — creaking ice, distant laughter, whispers.',
        '',
        'Comparable: Portal 2 (puzzles), Little Nightmares (atmosphere),',
        'What Remains of Edith Finch (narrative).',
        '',
        'Available: Interactive / video game rights worldwide.',
    ], 'has_cta': True},
])

# ===========================
# BEIN AND THE WORLD OF COLORS — ALL PITCHES
# ===========================
make_pitch(f"{outdir}/Bein_World_of_Colors_Master_Pitch.pdf", [
    {'title': 'BEIN AND THE WORLD OF THE COLORS', 'subtitle': 'An Emotional Fantasy by Eleanor Lian', 'body': [
        'Genre: Emotional Fantasy / Coming-of-Age / Magical Realism',
        'Format: Children\'s Novel / Illustrated Story',
        'Status: Complete manuscript, 23+ concept images, full character roster',
        '',
        'LOGLINE: A boy discovers that every emotion has a color — and that colors',
        'can heal, hurt, and reshape the world around him.',
        '',
        'STORY: Bein is a quiet boy who sees the world differently — literally.',
        'Every emotion around him manifests as a visible color. Anger is crimson.',
        'Joy is gold. Sadness is deep blue. Fear is grey. When his family faces',
        'a crisis that drains all color from his world, Bein must journey through',
        'the Land of Faded Hues to find the Source of All Colors and restore',
        'not just his world, but the emotional connections that hold his family',
        'together.',
        '',
        'THEMES: Emotional intelligence, family bonds, grief and healing,',
        'the beauty of feeling deeply.',
        '',
        'KEY CHARACTERS: Bein (protagonist), Mamma Caryl (nurturing mother),',
        'Sammy (protective older brother), Anton (curious friend),',
        'Fearfy (personification of fear), Fìonn (guide / personification of joy),',
        'Maith (the color witch), Wacky (chaos spirit).',
        '',
        'TONE: Inside Out meets The Little Prince meets The Boy and the Heron.',
        '',
        'RIGHTS AVAILABLE: Animation, Children\'s Book, Game, Merchandising.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Bein_World_of_Colors_Animation.pdf", [
    {'title': 'BEIN AND THE WORLD OF COLORS — Animated Film', 'subtitle': 'A feature animation for family audiences', 'body': [
        'Format: Animated Feature Film (80-90 min)',
        'Target: Family / Children 6+',
        'Tone: Pixar-level emotional depth meets Studio Ghibli visual beauty.',
        '',
        'PITCH: Bein\'s world offers a stunning visual concept — emotions as colors,',
        'the Land of Faded Hues as a greyscale wasteland, the Source of All Colors',
        'as a kaleidoscopic climax. Perfect for 2D or 3D animation.',
        '',
        'STRUCTURE:',
        '• Act 1: The Colorful World — Bein\'s gift, family bonds, the crisis hits',
        '• Act 2: The Fading — journey through the greyscale world, meeting',
        '  Fearfy, Maith, and the other personifications',
        '• Act 3: The Source — the truth about colors and feelings, the choice',
        '  that brings everything back, but changed',
        '',
        'EMOTIONAL BEAT: The film teaches children that ALL emotions are valid —',
        'sadness, fear, and anger have purpose alongside joy and love.',
        '',
        'Comparable: Inside Out, The Boy and the Heron, Song of the Sea.',
        '',
        'Available: Animation rights worldwide.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Bein_World_of_Colors_Series.pdf", [
    {'title': 'BEIN AND THE WORLD OF COLORS — Animated Series', 'subtitle': 'An episodic journey through emotions', 'body': [
        'Format: Animated Series, 12-24 episodes × 22 min',
        'Target: Children 4-10 and their families',
        'Platform: Netflix, Cartoon Network, Disney+, Nick Jr.',
        '',
        'PITCH: An episodic animated series where each episode explores a different',
        'color and the emotion it represents. Educational without being preachy.',
        '',
        'EPISODE STRUCTURE:',
        '• Ep 1-3: Introduction to color magic, meeting Fìonn',
        '• Ep 4-7: Exploration of primary emotions (red=anger, blue=sadness, yellow=joy)',
        '• Ep 8-12: Secondary emotions (green=jealousy, pink=love, grey=fear)',
        '• Ep 13-18: Maith\'s lessons — colors can mix and change',
        '• Ep 19-24: The Spectrum — all colors together, emotional balance',
        '',
        'EDUCATIONAL VALUE: Social-emotional learning (SEL) curriculum alignment.',
        'Each episode comes with discussion guides for parents and teachers.',
        '',
        'Comparable: Bluey, Hilda, Steven Universe (emotionally).',
        '',
        'Available: Series adaptation rights worldwide.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Bein_World_of_Colors_Book.pdf", [
    {'title': 'BEIN AND THE WORLD OF THE COLORS — Children\'s Books', 'subtitle': 'Illustrated book series', 'body': [
        'Format: Illustrated Children\'s Book Series',
        'Target: Ages 4-8 (picture books), Ages 6-10 (chapter books)',
        '',
        'PITCH: A series of beautifully illustrated books that teach emotional',
        'intelligence through color. Each book explores a different emotion/color',
        'with Bein and his friends.',
        '',
        'BOOK SERIES PLAN:',
        '• "Bein and the Red Storm" (Anger)',
        '• "Bein and the Blue River" (Sadness)',
        '• "Bein and the Golden Sun" (Joy)',
        '• "Bein and the Grey Fog" (Fear)',
        '• "Bein and the Green Vine" (Jealousy)',
        '• "Bein and the Pink Blossom" (Love)',
        '• "Bein and the Rainbow" (All emotions together)',
        '',
        'Each book: 32-48 pages, full-color illustrations, standard picture book format.',
        'Also available as: board books (toddler), activity books, coloring books.',
        '',
        'Comparable: The Color Monster, The Boy and the Colors, Ruby Finds a Worry.',
        '',
        'Available: Publishing rights worldwide (multiple languages).',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Bein_World_of_Colors_Game.pdf", [
    {'title': 'BEIN AND THE WORLD OF COLORS — Video Game', 'subtitle': 'Educational color-puzzle adventure', 'body': [
        'Format: Educational Puzzle-Adventure',
        'Platform: Mobile (iOS/Android), Nintendo Switch, PC',
        'Target: Children 4-10, educational market',
        '',
        'PITCH: A gentle puzzle game where players help Bein restore color to the',
        'world by solving emotion-based puzzles. Each puzzle teaches emotional',
        'recognition and regulation.',
        '',
        'GAME MECHANICS:',
        '• Color mixing puzzles — combine emotions to create nuanced feelings',
        '• Empathy challenges — recognize emotions in character faces/scenarios',
        '• Memory matching — color-emotion pair recall',
        '• Breathing exercises — interactive calming activities tied to colors',
        '• Free play mode — paint the world with colors, no pressure',
        '',
        'ACCESSIBILITY:',
        '• No text required (pre-reader friendly with voice guidance)',
        '• Colorblind-friendly mode (patterns + symbols alongside colors)',
        '• Parent dashboard with progress tracking',
        '',
        'Comparable: Toca Boca, Sago Mini, Endless Alphabet.',
        '',
        'Available: Educational game rights worldwide.',
    ], 'has_cta': True},
])

make_pitch(f"{outdir}/Bein_World_of_Colors_Merch.pdf", [
    {'title': 'BEIN AND THE WORLD OF COLORS — Merchandising', 'subtitle': 'Educational & playful products', 'body': [
        'Category: Toys, Apparel, Stationery, Educational Materials',
        'Target: Children 2-10, parents, educators',
        '',
        'PRODUCT LINES:',
        '',
        '1. Color Emotion Toys:',
        '• Plush: Bein, Fearfy (grey fuzzy), Fìonn (gold sparkle), Maith (color-changing)',
        '• Color mixing playset — physical color wheels for emotion identification',
        '• Emotion flashcards deck (52 cards, one per emotion+color)',
        '',
        '2. Apparel:',
        '• Color-changing t-shirts (reacts to body heat, shows Bein\'s world)',
        '• Mood socks — different emotion-color patterns',
        '• Pajama sets with celestial color patterns',
        '',
        '3. Educational:',
        '• Classroom kit: posters, worksheets, emotion wheel, teacher guide',
        '• Coloring book series (one per color/emotion)',
        '• App subscription: daily emotion check-in + color activities',
        '',
        '4. Home / Stationery:',
        '• "My Feelings" journal for kids',
        '• Mood lamp — changes color based on button press (emotion selector)',
        '• Bedtime story book + plush bundle sets',
        '',
        'Comparable: Inside Out merchandising line, The Color Monster products.',
        '',
        'Available: Global merchandising and licensing rights.',
    ], 'has_cta': True},
])

# Clean up any temp PPTX files
for f in os.listdir(outdir):
    if f.endswith('.pptx'):
        os.remove(os.path.join(outdir, f))

# Count results
pdfs = sorted([f for f in os.listdir(outdir) if f.endswith('.pdf')])
print(f"✅ Generated {len(pdfs)} pitch PDFs total\n")
for p in pdfs:
    sz = os.path.getsize(os.path.join(outdir, p))
    print(f"   {p} ({sz/1024:.0f} KB)")
