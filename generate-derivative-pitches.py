#!/usr/bin/env python3
"""Generate individual pitch deck PDFs for each derivative of Sigillum Maximum."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

GOLD = RGBColor(0xD4, 0xA7, 0x47)
DARK = RGBColor(0x0A, 0x0A, 0x0F)
CREAM = RGBColor(0xF0, 0xE8, 0xD8)
LIGHT = RGBColor(0xE0, 0xDC, 0xD5)
MUTED = RGBColor(0xA0, 0x98, 0x88)

def set_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gold_bar(slide, top=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, prs.slide_width, Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def add_title(slide, text, top=0.5):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.4), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = CREAM
    p.font.bold = True
    p.font.name = 'Calibri'
    
def add_subtitle(slide, text, top=1.1):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = MUTED
    p.font.name = 'Calibri'

def add_bullet(slide, items, left=0.8, top=1.7, width=8.4, font_size=13, color=LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox

def save_pitch(path, slides):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        s = data
        set_bg(slide)
        add_gold_bar(slide)
        
        add_title(slide, s['title'])
        add_subtitle(slide, s.get('subtitle', ''))
        
        if 'bullets' in s:
            add_bullet(slide, s['bullets'])
        
        if s.get('has_cta'):
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(8.4), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "Interested? Contact: info@robertocalvoproductions.com"
            p.font.size = Pt(12)
            p.font.color.rgb = GOLD
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
    
    prs.save(path)
    return os.path.getsize(path)

# Temp directory for individual derivatives
output_dir = "/root/.openclaw/workspace/sigillum-maximum-site/derivative-pitches"
os.makedirs(output_dir, exist_ok=True)

# ==================== 1. ANIME SERIES PITCH ====================
save_pitch(f"{output_dir}/Sigillum_Maximum_Anime_Series.pdf", [
    {
        'title': 'Sigillum Maximum — Anime Series Pitch',
        'subtitle': 'A princess. Seven seals. One choice that will save — or shatter — the Great Universal Realm.',
        'bullets': [
            'Format: TV Anime Series, 12–24 episodes × 24 minutes',
            'Genre: Epic Fantasy / Coming-of-Age / Adventure',
            'Target: Teens & Young Adults (13+), global audience',
            'Language: English (dub-ready for JP, FR, ES markets)',
            '',
            "Story Arc Structure:",
            '• Act 1 (Ep 1-4): Nuvolandia shattered — Aileen enters the Elf Kingdom',
            '• Act 2 (Ep 5-12): Journey through Gnome, Merfolk, Dwarf Kingdoms',
            '• Act 3 (Ep 13-20): Storm Kingdom, betrayal, the truth about Dorcha',
            '• Finale (Ep 21-24): Realm of the Great Light — the final secret',
            '',
            'Visual Style: Dark fantasy with vibrant magical realms, painters CGI backgrounds + 2D character animation.',
            'Tonal Reference: The Ancient Magus\' Bride meets Made in Abyss meets Frieren.',
            '',
            'Co-Production Opportunity: Seeking Japanese/Western studio partnership for simulcast release.',
            'Available: Worldwide streaming rights (excluding UK publishing home market).',
        ],
        'has_cta': True,
    }
])

# ==================== 2. FILM TRILOGY PITCH ====================
save_pitch(f"{output_dir}/Sigillum_Maximum_Film_Trilogy.pdf", [
    {
        'title': 'Sigillum Maximum — Film Trilogy Pitch',
        'subtitle': 'Three films. Seven realms. One epic cinematic journey.',
        'bullets': [
            'Format: Film Trilogy (Live-Action or CGI Animation)',
            'Runtime: Film 1 — 120min / Film 2 — 130min / Film 3 — 140min',
            'Target: Family / Fantasy / Adventure audiences (PG-13 / 12A)',
            '',
            'Film 1 — The Shattering (120 min)',
            '• Nuvolandia invaded by the Black Knight. King Eadwig captured.',
            '• Queen Eloria turned to stone. Aileen escapes through the hidden passage.',
            '• Arrival at the Elf Kingdom. Meeting Aeltiàfisar, the Knight of the Golden Light.',
            '• First secrets of the Sigillum revealed. Quest begins.',
            '',
            'Film 2 — The Journey (130 min)',
            '• Traversal of five kingdoms: Gnome, Merfolk, Dwarf, Storm, the abandoned lands.',
            '• Grogher joins the party. The dark truth about Dorcha\'s origins surfaces.',
            '• Majory\'s curse revealed. Alliance with the orcs.',
            '• Midpoint twist: the Sigillum demands a sacrifice.',
            '',
            'Film 3 — The Final Seal (140 min)',
            '• Race to the Realm of the Great Light.',
            '• The Black Knight\'s true identity. Dorcha\'s ultimate choice.',
            '• The seventh secret: what it really means to break all seals.',
            '• Final battle. Resolution. The cost of saving the universe.',
            '',
            'Visual Style: Rich CGI akin to Ghibli meets Final Fantasy.',
            'Comparable: The NeverEnding Story, How to Train Your Dragon, The Dark Crystal.',
            '',
            'Available: Film rights worldwide.',
        ],
        'has_cta': True,
    }
])

# ==================== 3. VIDEO GAME PITCH ====================
save_pitch(f"{output_dir}/Sigillum_Maximum_Video_Game.pdf", [
    {
        'title': 'Sigillum Maximum — Video Game Pitch',
        'subtitle': 'An open-world Action RPG across the Seven Realms.',
        'bullets': [
            'Genre: Action RPG / Open World / Story-Driven',
            'Platform: PC, PlayStation, Xbox, Nintendo Switch',
            'Target: Fantasy RPG fans (Elden Ring, Zelda, Genshin Impact audiences)',
            '',
            'Core Concept:',
            '• Play as Aileen, exploring all Seven Realms in a non-linear open world.',
            '• Each kingdom has unique biomes, enemies, puzzles, and magic systems.',
            '• Companion system: recruit Grogher, Dorcha, and others with distinct combat styles.',
            '• Magic system based on the Seven Secrets — unlock new abilities per realm.',
            '',
            'Key Features:',
            '• 7 massive regions, each ~20-30 sq km of explorable terrain',
            '• Mount system: Raertha (flight), Sidae (ground/fast travel), sea mounts',
            '• Dynamic day/night cycle with realm-specific weather',
            '• Crafting, alchemy, and rune-enchantment mechanics',
            '• Morality system tied to Dorcha\'s dark/light alignment',
            '• Voice-acted main quest (30+ hours) + hundreds of side quests',
            '',
            'Market Fit:',
            '• Strong demand for narrative-driven fantasy RPGs ($28B+ market)',
            '• Anime art style resonates with Genshin/Star Rail player base',
            '• Unproven IP with built-in lore depth for expansion',
            '• Potential for live-service elements (realm events, seasonal content)',
            '',
            'Opportunity: Co-development with AA/indie RPG studio.',
            'Available: Interactive media / video game rights.',
        ],
        'has_cta': True,
    }
])

# ==================== 4. MERCHANDISING PITCH ====================
save_pitch(f"{output_dir}/Sigillum_Maximum_Merchandising.pdf", [
    {
        'title': 'Sigillum Maximum — Merchandising & Licensing Pitch',
        'subtitle': 'A rich fantasy universe built for consumer products.',
        'bullets': [
            'Category Overview:',
            '',
            'Collectibles & Figures:',
            '• Premium 7" action figures (Aileen, Dorcha, Grogher, Aeltiàfisar)',
            '• Realm-themed statue collection (7 pieces, one per kingdom)',
            '• Plush line: Raertha the unicorn, Sidae the winged lion, Hercules',
            '• Blind box mini-figures (12 characters, 2 rare variants)',
            '',
            'Apparel & Accessories:',
            '• T-shirt line (character art, realm crests, Sigillum symbols)',
            '• Hoodies with realm-color schemes and runic embroidery',
            '• Jewelry: Sigillum pendant, realm rings set, Aileen\'s moonstone crown replica',
            '• Cosplay-grade costumes (Aileen dress, Dorcha armor, Grogher orc gear)',
            '',
            'Stationery & Home:',
            '• Art book: "The Realms of Sigillum Maximum" (200+ pages of concept art)',
            '• Realm-themed journal set (7 journals with unique covers)',
            '• Posters, mouse pads, phone cases, enamel pins',
            '• Board game adaptation: "Sigillum — The Quest for the Seven Secrets"',
            '',
            'Digital & Publishing:',
            '• Manga/comic adaptation (prequels, side stories)',
            '• Art print collection (signed, limited edition)',
            '• Soundtrack album (digital + vinyl)',
            '• Cookbook: "Feasts of the Seven Realms"',
            '',
            'Target Partners: Licensing agencies, toy manufacturers, publishers.',
            'Available: Global merchandising and licensing rights.',
        ],
        'has_cta': True,
    }
])

# ==================== 5. MANGA / COMIC PITCH ====================
save_pitch(f"{output_dir}/Sigillum_Maximum_Manga.pdf", [
    {
        'title': 'Sigillum Maximum — Manga / Comic Adaptation Pitch',
        'subtitle': 'Expand the universe through sequential art.',
        'bullets': [
            'Proposed Format: Monthly serialized manga (or Western comic)',
            'Target: Manga/anime readers, fantasy comic fans',
            'Demographic: Shônen / Seinen crossover (Teens + Young Adults)',
            '',
            'Pitch: The Sigillum Maximum world is a natural fit for manga. Its visual richness, kingdom-based story structure, and character-driven arcs translate perfectly to panel-by-panel storytelling.',
            '',
            'Potential Titles:',
            '• "Sigillum Maximum: The Lost Realm" (prequel — the war before the story)',
            '• "Sigillum Maximum: Dorcha" (spin-off exploring the Black Knight\'s origin)',
            '• "Sigillum Maximum: Side Stories" (anthology — glimpses into each realm)',
            '• "Sigillum Maximum: The Art of the Seven Secrets" (art book / guide)',
            '',
            'Format Details:',
            '• Volume 1: ~200 pages, $12.99 digital/$24.99 print',
            '• Planned: 5-7 volumes for complete prequel arc',
            '• Digital-first (comiXology, Manga Plus, Webtoons) + print collection',
            '• Bilingual release: Italian (original) + English (global)',
            '',
            'Art Style: Anime/manga aesthetic consistent with animated adaptation.',
            'Comparable to: The Ancient Magus\' Bride, Witch Hat Atelier, Delicious in Dungeon.',
            '',
            'Available: Comic publishing rights worldwide.',
        ],
        'has_cta': True,
    }
])

# ==================== 6. BOARD GAME PITCH ====================
save_pitch(f"{output_dir}/Sigillum_Maximum_Board_Game.pdf", [
    {
        'title': 'Sigillum Maximum — Board Game Pitch',
        'subtitle': 'A competitive & cooperative journey across the Seven Realms.',
        'bullets': [
            'Format: Board game (2-6 players, 60-120 min, ages 12+)',
            'Genre: Adventure / Strategy / Cooperative',
            'Target: Fantasy board gamers (Gloomhaven, Mage Knight, Nemesis)',
            '',
            'Core Concept: Players are heroes racing across the Seven Realms to collect the secrets of the Sigillum before the Black Knight assembles his army.',
            '',
            'Game Mechanics:',
            '• Modular board: 7 realm tiles placed randomly each game',
            '• Each realm has unique challenges, allies, and secret cards',
            '• Character classes: Knight (Aeltiàfisar), Mage (Majory), Warrior (Grogher), Rogue (Dorcha), Healer (Aileen)',
            '• Cooperative play with competitive scoring (who contributed most)',
            '• Resource management: light essence, realm tokens, alliance markers',
            '• Event deck: story cards that unfold the narrative mid-game',
            '',
            'Expansion Potential:',
            '• "Nuvolandia Rising" (exploration expansion)',
            '• "The Black Knight\'s Army" (vs. mode expansion)',
            '• "Dorcha\'s Dilemma" (story expansion with branching paths)',
            '• Miniature upgrade pack (premium painted figures)',
            '',
            'Comparable to: Gloomhaven meets The Lord of the Rings board game.',
            '',
            'Opportunity: Licensing to Asmodee, CMON, or indie board game publisher.',
            'Available: Tabletop game rights worldwide.',
        ],
        'has_cta': True,
    }
])

# ==================== 7. THEME PARK / EXPERIENCE PITCH ====================
save_pitch(f"{output_dir}/Sigillum_Maximum_Experiences.pdf", [
    {
        'title': 'Sigillum Maximum — Immersive Experiences Pitch',
        'subtitle': 'Bring the Seven Realms to life.',
        'bullets': [
            'Concept: Immersive entertainment and location-based experiences.',
            'Target: Theme parks, pop-up attractions, escape rooms.',
            '',
            'Experience Concepts:',
            '',
            '1. "The Hidden Passage" — Escape Room Experience:',
            '• Players navigate a series of rooms themed after realms',
            '• Solve puzzles to unlock the secrets of the Sigillum',
            '• 45-60 min experience, 4-8 players',
            '• Available as permanent installation or touring pop-up',
            '',
            '2. "Flight Across the Realms" — VR / Motion Simulator:',
            '• 4D motion ride: fly on Raertha across all seven kingdoms',
            '• Wind, heat, water effects synchronized to each realm biome',
            '• Can be installed in theme parks, malls, VR arcades',
            '',
            '3. "Sigillum: Live" — Immersive Theater:',
            '• Live-action interactive theater with professional actors',
            '• Audience becomes part of Aileen\'s quest',
            '• Deployable as seasonal event or permanent show',
            '',
            '4. Pop-Up Store: "Merchant of the Seven Realms":',
            '• Themed retail experience selling merchandising line',
            '• Photo-op zones with life-size character statues',
            '• Perfect for anime conventions, comic cons, mall activations',
            '',
            'Comparable: Wizarding World, Stranger Things experience, Avatar: Flight of Passage.',
            '',
            'Available: Location-based entertainment rights worldwide.',
        ],
        'has_cta': True,
    }
])

# Clean up temp PPTX
for f in os.listdir(output_dir):
    if f.endswith('.pptx'):
        os.remove(os.path.join(output_dir, f))

print(f"✅ Generated {len([f for f in os.listdir(output_dir) if f.endswith('.pdf')])} derivative pitch PDFs")
for f in sorted(os.listdir(output_dir)):
    size = os.path.getsize(os.path.join(output_dir, f))
    print(f"   {f} ({size/1024:.0f} KB)")
