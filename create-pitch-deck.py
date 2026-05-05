#!/usr/bin/env python3
"""Create a professional pitch deck PDF for Sigillum Maximum."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
GOLD = RGBColor(0xD4, 0xA7, 0x47)
DARK = RGBColor(0x0A, 0x0A, 0x0F)
DARK2 = RGBColor(0x14, 0x14, 0x1E)
CREAM = RGBColor(0xF0, 0xE8, 0xD8)
LIGHT = RGBColor(0xE0, 0xDC, 0xD5)
MUTED = RGBColor(0xA0, 0x98, 0x88)
DARK_GOLD = RGBColor(0x8B, 0x6F, 0x30)

def set_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gold_bar(slide, top=0, height=Inches(0.04)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), top, prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def add_text_box(slide, text, left, top, width, height, font_size=18, color=LIGHT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, items, left, top, width, height, font_size=16, color=MUTED):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return txBox

def add_quote_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x1A)
    shape.line.color.rgb = RGBColor(0xD4, 0xA7, 0x47)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = f'"{text}"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GOLD
    p.font.name = 'Calibri'
    return shape

# ============== SLIDE 1: TITLE ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_gold_bar(slide, Inches(2.8), Inches(0.03))

add_text_box(slide, "ROBERTO CALVO PRODUCTIONS LTD", 
    Inches(1), Inches(1.2), Inches(11), Inches(0.6),
    font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "SIGILLUM MAXIMUM",
    Inches(1), Inches(3.0), Inches(11), Inches(1.2),
    font_size=60, color=CREAM, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "The Seal of the Seven Secrets",
    Inches(1), Inches(4.0), Inches(11), Inches(0.8),
    font_size=28, color=GOLD, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "An Epic Fantasy Universe Ready for Adaptation",
    Inches(1), Inches(4.8), Inches(11), Inches(0.6),
    font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "by Eleanor Lian · Published by Roberto Calvo Productions Ltd",
    Inches(1), Inches(5.8), Inches(11), Inches(0.5),
    font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# ============== SLIDE 2: LOGLINE ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "The Story",
    Inches(1), Inches(0.5), Inches(5), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

add_text_box(slide, "LOGLINE",
    Inches(1), Inches(1.1), Inches(5), Inches(0.4),
    font_size=12, color=GOLD, bold=True)

add_quote_box(slide, "A princess. Seven seals. One choice that will save — or shatter — the Great Universal Realm.",
    Inches(1), Inches(1.8), Inches(11), Inches(1.2))

add_text_box(slide, "SYNOPSIS",
    Inches(1), Inches(3.3), Inches(5), Inches(0.4),
    font_size=12, color=GOLD, bold=True)

add_bullet_text(slide, [
    "Princess Aileen is 13 when a mysterious Black Knight invades Nuvolandia, turning her mother to stone and capturing her father.",
    "The Sigillum Maximum — the ancient golden scroll protecting all seven realms — has been stolen.",
    "Sent through a hidden passage to the Elf Kingdom, Aileen must find Aeltiàfisar, the legendary Knight of the Golden Light.",
    "Together with fierce orc Grogher and the enigmatic Dorcha, she embarks on a perilous quest across all seven realms.",
    "Each kingdom holds a secret of the Sigillum — but the final secret may demand a sacrifice no one is ready to make."
], Inches(1), Inches(3.8), Inches(11), Inches(3.2), font_size=14, color=LIGHT)

# ============== SLIDE 3: THE WORLD ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "The Great Universal Realm",
    Inches(1), Inches(0.5), Inches(5), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

add_text_box(slide, "A universe of seven kingdoms, each with its own magic, culture, and secret.",
    Inches(1), Inches(1.2), Inches(11), Inches(0.5),
    font_size=16, color=MUTED)

realms = [
    ("☁️", "Nuvolandia", "The Cloud Kingdom. Moonstone castles, beings in telepathic contact with the skies."),
    ("🌿", "Elf Kingdom", "Towering forests of luminous trees. Home to Aeltiàfisar and King Baelkers."),
    ("⛰️", "Gnome Kingdom", "Kingdom of the Two Rainbows. Subterranean crystals, Pallafiocco tournaments."),
    ("🌊", "Merfolk Kingdom", "Coral palaces, mother-of-pearl gowns. Secrets of the deep unknown to elves."),
    ("🏔️", "Dwarf Kingdom", "Forged in mountains. Masters of craft, stone, and ancient metallurgy."),
    ("⚡", "Storm Kingdom", "Perpetual tempests. Home to orcs. A secret buried in thunder."),
    ("✨", "Realm of the Great Light", "Legendary domain of Immortal Sovereigns. Source of the Sigillum Maximum."),
]

y = 2.0
for icon, name, desc in realms:
    add_text_box(slide, f"{icon} {name}",
        Inches(1), Inches(y), Inches(3.5), Inches(0.4),
        font_size=15, color=GOLD, bold=True)
    add_text_box(slide, desc,
        Inches(4.5), Inches(y), Inches(7.5), Inches(0.4),
        font_size=13, color=MUTED)
    y += 0.65

# ============== SLIDE 4: CHARACTERS ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "Main Characters",
    Inches(1), Inches(0.5), Inches(5), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

chars = [
    ("Aileen", "Protagonist", "Princess of Nuvolandia. 13 years old. Brave, compassionate, growing into a leader."),
    ("Dorcha", "The Knight with Darkness in His Blood", "Raised by orcs. Carries a darkness that may save — or doom — the realm."),
    ("Grogher", "The Fierce Orc", "Separated from his family. Immense strength and loyalty. Accompanied by winged lion Sidae."),
    ("Aeltiàfisar", "Knight of the Golden Light", "Legendary elven knight. Ancient, wise. Aileen's guide and mentor."),
    ("Majory", "Gnome Princess", "Competitive, caring. Must break an ancient curse upon her family."),
    ("King Eadwig", "King of Nuvolandia", "Turned to stone. His fate drives Aileen's quest."),
    ("Queen Eloria", "Queen of Nuvolandia", "Transformed into a statue. Symbol of all Aileen fights to restore."),
]

y = 1.5
for name, role, desc in chars:
    add_text_box(slide, f"{name}",
        Inches(1), Inches(y), Inches(2.5), Inches(0.35),
        font_size=15, color=CREAM, bold=True)
    add_text_box(slide, role,
        Inches(3.5), Inches(y), Inches(3.5), Inches(0.35),
        font_size=11, color=GOLD)
    add_text_box(slide, desc,
        Inches(7), Inches(y), Inches(5.5), Inches(0.35),
        font_size=12, color=MUTED)
    y += 0.6

# ============== SLIDE 5: THEMES ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "Themes & Core Conflict",
    Inches(1), Inches(0.5), Inches(11), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

add_text_box(slide, "A story about growing up, loss, love, and the choices that define who we become.",
    Inches(1), Inches(1.2), Inches(11), Inches(0.5),
    font_size=16, color=MUTED)

themes = [
    ("Courage in Vulnerability", "Aileen is not a warrior — she is a girl who must find strength in kindness, empathy, and the will to keep going."),
    ("The Cost of Power", "Every secret of the Sigillum has a price. The final seal demands a sacrifice that tests the very meaning of love."),
    ("Found Family", "Aileen, Grogher, Dorcha — each broken in their own way — form a bond stronger than blood."),
    ("Light vs. Darkness", "Not a simple battle. Dorcha embodies both. The story asks: can darkness be redeemed?"),
    ("Loss & Transformation", "Every character loses something. The question is what they become through that loss."),
]

y = 2.0
for title, desc in themes:
    add_text_box(slide, f"✦ {title}",
        Inches(1), Inches(y), Inches(4), Inches(0.35),
        font_size=16, color=GOLD, bold=True)
    add_text_box(slide, desc,
        Inches(5), Inches(y), Inches(7), Inches(0.35),
        font_size=13, color=MUTED)
    y += 0.85

# ============== SLIDE 6: MARKET ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "Market & Audience",
    Inches(1), Inches(0.5), Inches(11), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

# Left column
add_text_box(slide, "Target Demographics",
    Inches(1), Inches(1.5), Inches(5), Inches(0.4),
    font_size=18, color=GOLD, bold=True)

add_bullet_text(slide, [
    "Teens & Young Adults (13–25) — core demographic",
    "Fantasy readers (Sanderson, Maas, Tolkien audiences)",
    "Anime & manga fans worldwide",
    "Western animation / streaming audiences",
    "Family co-viewing appeal (emotional depth, not just action)"
], Inches(1), Inches(2.0), Inches(5), Inches(3), font_size=14, color=LIGHT)

# Right column
add_text_box(slide, "Market Trends",
    Inches(7), Inches(1.5), Inches(5), Inches(0.4),
    font_size=18, color=GOLD, bold=True)

add_bullet_text(slide, [
    "Fantasy adaptation is the dominant genre in streaming (2024-26)",
    "Strong demand for female-led epic fantasy (Arcane, The Dragon Prince)",
    "Anime-style Western co-productions growing rapidly",
    "Multi-platform potential: series + games + publishing",
    "Comparable IP valuations: $50M–$200M for top fantasy adaptations"
], Inches(7), Inches(2.0), Inches(5), Inches(3), font_size=14, color=LIGHT)

add_text_box(slide, "The fantasy adaptation market is projected to grow 12.4% annually through 2030.",
    Inches(1), Inches(5.8), Inches(11), Inches(0.5),
    font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# ============== SLIDE 7: ADAPTATION ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "Proposed Format",
    Inches(1), Inches(0.5), Inches(11), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

add_text_box(slide, "A project designed for both anime series and film trilogies — ready for studio interest.",
    Inches(1), Inches(1.2), Inches(11), Inches(0.5),
    font_size=16, color=MUTED)

# Two columns
add_text_box(slide, "🎬 Option A: Anime TV Series",
    Inches(1), Inches(2.0), Inches(5), Inches(0.4),
    font_size=22, color=GOLD, bold=True)

add_bullet_text(slide, [
    "12–24 episodes, 24 minutes each",
    "One major kingdom arc per 3-4 episodes",
    "Season 1: Nuvolandia → Elf Kingdom → Gnome Kingdom",
    "Season 2: Merfolk → Dwarf → Storm → Great Light",
    "Open for simulcast / co-production with Japanese studios",
    "Comparable: The Ancient Magus' Bride, Made in Abyss, Frieren"
], Inches(1), Inches(2.6), Inches(5), Inches(3.5), font_size=14, color=LIGHT)

add_text_box(slide, "🎥 Option B: Live-Action / CGI Trilogy",
    Inches(7), Inches(2.0), Inches(5), Inches(0.4),
    font_size=22, color=GOLD, bold=True)

add_bullet_text(slide, [
    "Film 1: The Shattering — 120 min",
    "Film 2: The Journey — 130 min",
    "Film 3: The Final Seal — 140 min",
    "Visual style: rich CGI akin to Ghibli meets Final Fantasy",
    "Also viable as high-end limited series (8-10 episodes)",
    "Comparable: The NeverEnding Story, How to Train Your Dragon"
], Inches(7), Inches(2.6), Inches(5), Inches(3.5), font_size=14, color=LIGHT)

# ============== SLIDE 8: RIGHTS & ASSETS ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "What We Offer",
    Inches(1), Inches(0.5), Inches(11), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

add_text_box(slide, "Rights Available:",
    Inches(1), Inches(1.5), Inches(5), Inches(0.4),
    font_size=18, color=GOLD, bold=True)

add_bullet_text(slide, [
    "Anime / TV / Streaming adaptation rights",
    "Film adaptation rights (live-action or animated)",
    "Video game / interactive media rights",
    "Merchandising & licensing rights",
    "Publishing rights (expanded universe, art books)"
], Inches(1), Inches(2.0), Inches(5), Inches(3), font_size=14, color=LIGHT)

add_text_box(slide, "Existing Assets:",
    Inches(7), Inches(1.5), Inches(5), Inches(0.4),
    font_size=18, color=GOLD, bold=True)

add_bullet_text(slide, [
    "Full manuscript — complete, ready for production",
    "55+ character & concept artworks (AI-assisted, iterable)",
    "Fully developed world bible (7 realms, history, magic system)",
    "Detailed character profiles and arc outlines",
    "Proven cross-platform model (BlaBlaSell, FanCake ecosystem)"
], Inches(7), Inches(2.0), Inches(5), Inches(3), font_size=14, color=LIGHT)

# ============== SLIDE 9: ROADMAP ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide)

add_text_box(slide, "Roadmap",
    Inches(1), Inches(0.5), Inches(11), Inches(0.7),
    font_size=36, color=CREAM, bold=True)

add_text_box(slide, "Phase 1 — Now",
    Inches(1), Inches(1.5), Inches(4), Inches(0.4),
    font_size=18, color=GOLD, bold=True)
add_bullet_text(slide, [
    "✦ Landing page live (sigillummaximum.com)",
    "✦ Pitch deck & production materials ready",
    "✦ Securing initial producer conversations",
    "✦ Finalizing manuscript for publisher submission"
], Inches(1), Inches(2.0), Inches(4), Inches(2.5), font_size=14, color=LIGHT)

add_text_box(slide, "Phase 2 — Q3 2026",
    Inches(5.5), Inches(1.5), Inches(4), Inches(0.4),
    font_size=18, color=GOLD, bold=True)
add_bullet_text(slide, [
    "✦ Publisher submission (US & UK markets)",
    "✦ Studio pitch meetings (anime & live-action)",
    "✦ Production partner identification",
    "✦ Expanded universe planning"
], Inches(5.5), Inches(2.0), Inches(4), Inches(2.5), font_size=14, color=LIGHT)

add_text_box(slide, "Phase 3 — 2027",
    Inches(10), Inches(1.5), Inches(4), Inches(0.4),
    font_size=18, color=GOLD, bold=True)
add_bullet_text(slide, [
    "✦ Co-production agreement",
    "✦ Pre-production & concept development",
    "✦ Casting / studio partnership",
    "✦ B2B launch at major industry events"
], Inches(10), Inches(2.0), Inches(4), Inches(2.5), font_size=14, color=LIGHT)

# Bottom
add_text_box(slide, "We are seeking: Co-production partners · Streamers · Publishers · Game developers",
    Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)

# ============== SLIDE 10: CONTACT ==============
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gold_bar(slide, Inches(2.8), Inches(0.03))

add_text_box(slide, "Let's Build Something Extraordinary",
    Inches(1), Inches(2.0), Inches(11), Inches(0.8),
    font_size=40, color=CREAM, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "Roberto Calvo Productions Ltd",
    Inches(1), Inches(3.2), Inches(11), Inches(0.6),
    font_size=20, color=GOLD, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "71-75 Shelton Street · London · WC2H 9JQ · United Kingdom",
    Inches(1), Inches(3.9), Inches(11), Inches(0.5),
    font_size=15, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "info@robertocalvoproductions.com",
    Inches(1), Inches(4.5), Inches(11), Inches(0.5),
    font_size=18, color=GOLD, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "robertocalvoproductions.com · sigillummaximum.com",
    Inches(1), Inches(5.1), Inches(11), Inches(0.5),
    font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "Eleanor Lian presents: Sigillum Maximum — The Seal of the Seven Secrets",
    Inches(1), Inches(6.2), Inches(11), Inches(0.4),
    font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/root/.openclaw/workspace/sigillum-maximum-site/Sigillum_Maximum_Pitch_Deck.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
print(f"   File size: {os.path.getsize(output_path) / 1024:.0f} KB")
print(f"   Slides: {len(prs.slides)}")
